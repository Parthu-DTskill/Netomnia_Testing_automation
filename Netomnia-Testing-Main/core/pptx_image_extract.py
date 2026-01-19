import os
import io
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

def extract_images_from_pptx(pptx_path, output_dir, prefix="pptx_image"):
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(pptx_path)

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    extracted_files = []
    seen_hashes = set()
    counter = 1

    for slide_index, slide in enumerate(prs.slides, 1):

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_bytes = shape.image.blob
            image_hash = hashlib.md5(image_bytes).hexdigest()

            if image_hash in seen_hashes:
                continue
            seen_hashes.add(image_hash)

            ext = shape.image.ext.lower()
            filename = f"{prefix}_{slide_index}_{counter}.{ext}"
            output_path = os.path.join(output_dir, filename)

            with Image.open(io.BytesIO(image_bytes)) as img:
                img.save(output_path)

            extracted_files.append(output_path)
            counter += 1

    return extracted_files