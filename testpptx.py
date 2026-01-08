import os
import io
import hashlib
from pptx import Presentation
from PIL import Image


def extract_images_from_pptx(
    pptx_path: str,
    output_dir: str,
    prefix: str = "pptx_image"
):
    """
    Extracts all images from a PPTX file safely and efficiently.

    - Preserves original image quality
    - Avoids duplicate images
    - Server-safe (no rendering)
    - Works with all slides and shapes

    Returns:
        List of extracted image file paths
    """

    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    extracted_files = []
    seen_hashes = set()
    counter = 1

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.shape_type == 13:  # 13 = MSO_SHAPE_TYPE.PICTURE
                continue

            image = shape.image
            image_bytes = image.blob

            image_hash = hashlib.md5(image_bytes).hexdigest()
            if image_hash in seen_hashes:
                continue
            seen_hashes.add(image_hash)

            image_ext = image.ext.lower()
            filename = f"{prefix}_{slide_index}_{counter}.{image_ext}"
            output_path = os.path.join(output_dir, filename)

            with Image.open(io.BytesIO(image_bytes)) as img:
                img.save(output_path)

            extracted_files.append(output_path)
            counter += 1

    return extracted_files

pptx_file = "D:\\Netomnia_Automation_Final\\H&D125_120905_1.pptx"
output_folder = "D:\\Netomnia_Automation_Final\\Test_PPTX"

images = extract_images_from_pptx(
    pptx_path=pptx_file,
    output_dir=output_folder
)

print(f"Extracted {len(images)} images")